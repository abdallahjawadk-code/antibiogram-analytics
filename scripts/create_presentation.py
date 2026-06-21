"""
AntibioGram Pro — Official Presentation Generator
Generates a professional bilingual (Arabic/English) PPTX presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree
import copy, os

# ── Color Palette ─────────────────────────────────────────────────────────────
TEAL      = RGBColor(0x0d, 0x94, 0x88)   # teal-600
CYAN      = RGBColor(0x06, 0x82, 0xb2)   # cyan-600  (was 0891b2, fixed hex)
NAVY      = RGBColor(0x0f, 0x17, 0x2a)   # slate-900
DARK      = RGBColor(0x1e, 0x29, 0x3b)   # slate-800
MID       = RGBColor(0x47, 0x55, 0x69)   # slate-600
LIGHT_BG  = RGBColor(0xf0, 0xfd, 0xfa)   # teal-50
WHITE     = RGBColor(0xff, 0xff, 0xff)
VIOLET    = RGBColor(0x7c, 0x3a, 0xed)   # violet-600
ROSE      = RGBColor(0xe1, 0x1d, 0x48)   # rose-600
AMBER     = RGBColor(0xd9, 0x77, 0x06)   # amber-600
EMERALD   = RGBColor(0x05, 0x96, 0x69)   # emerald-600

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=0, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER, rtl=False, font_name="Calibri",
                italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    # RTL paragraph direction
    if rtl:
        pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1')
        run.font.name = "Traditional Arabic" if font_name == "Calibri" else font_name
    return txb


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK,
                  align=PP_ALIGN.LEFT, rtl=False, space_before=0, font_name="Calibri"):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = font_name if not rtl else "Traditional Arabic"
    if rtl:
        pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1')
    return p


def gradient_bg(slide, color1=NAVY, color2=TEAL):
    """Simulate gradient with two overlapping rectangles."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=color1)
    # A wide translucent strip for the gradient feel
    strip = slide.shapes.add_shape(1, Inches(0), Inches(4), Inches(13.33), Inches(3.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color2
    strip.line.fill.background()
    # Darken by blending — achieved via XML opacity
    spPr = strip._element.spPr
    solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is None:
            srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgb.set('val', '%02x%02x%02x' % (color2.red, color2.green, color2.blue))
        alpha_el = etree.SubElement(srgb, qn('a:alpha'))
        alpha_el.set('val', '70000')  # 70% opacity


def teal_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x02, 0x2c, 0x2c))


def add_section_header(slide, ar_text, en_text, y=0.18):
    add_rect(slide, 0, 0, 13.33, y + 0.72, fill_rgb=TEAL)
    add_textbox(slide, ar_text, 0.3, y - 0.05, 8, 0.55,
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, rtl=True)
    add_textbox(slide, en_text, 0.3, y - 0.05, 12.7, 0.55,
                font_size=13, bold=False, color=RGBColor(0xcc, 0xff, 0xf7),
                align=PP_ALIGN.LEFT)


def add_bullet_box(slide, title_ar, items_ar, x, y, w, h,
                   bg=LIGHT_BG, accent=TEAL):
    add_rect(slide, x, y, w, h, fill_rgb=bg,
             line_rgb=accent, line_width=1.5)
    # accent left bar
    add_rect(slide, x, y, 0.06, h, fill_rgb=accent)
    txb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.25), Inches(h - 0.2))
    tf = txb.text_frame
    tf.word_wrap = True
    # Title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = title_ar
    run.font.size      = Pt(14)
    run.font.bold      = True
    run.font.color.rgb = accent
    run.font.name      = "Traditional Arabic"
    pPr = p._p.get_or_add_pPr()
    pPr.set(qn('a:rtl'), '1')
    # Items
    for item in items_ar:
        add_paragraph(tf, item, font_size=12, color=DARK,
                      align=PP_ALIGN.RIGHT, rtl=True, space_before=3)


def add_formula_box(slide, label, formula, result, x, y, w=4.0, bg=RGBColor(0xec, 0xfd, 0xfb)):
    add_rect(slide, x, y, w, 1.15, fill_rgb=bg,
             line_rgb=TEAL, line_width=1.2)
    add_textbox(slide, label, x + 0.1, y + 0.05, w - 0.2, 0.3,
                font_size=11, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_textbox(slide, formula, x + 0.1, y + 0.3, w - 0.2, 0.42,
                font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                font_name="Cambria Math")
    add_textbox(slide, result, x + 0.1, y + 0.78, w - 0.2, 0.3,
                font_size=10, color=MID, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

gradient_bg(slide)

# Decorative circle top-right
c1 = slide.shapes.add_shape(9, Inches(10.5), Inches(-1.2), Inches(4), Inches(4))
c1.fill.solid(); c1.fill.fore_color.rgb = TEAL
spPr = c1._element.spPr
sf = spPr.find('.//' + qn('a:solidFill'))
if sf is not None:
    s = sf.find(qn('a:srgbClr'))
    if s is None:
        s = etree.SubElement(sf, qn('a:srgbClr')); s.set('val', '0d9488')
    a = etree.SubElement(s, qn('a:alpha')); a.set('val', '25000')
c1.line.fill.background()

# Decorative circle bottom-left
c2 = slide.shapes.add_shape(9, Inches(-0.8), Inches(5.8), Inches(3), Inches(3))
c2.fill.solid(); c2.fill.fore_color.rgb = CYAN
spPr2 = c2._element.spPr
sf2 = spPr2.find('.//' + qn('a:solidFill'))
if sf2 is not None:
    s2 = sf2.find(qn('a:srgbClr'))
    if s2 is None:
        s2 = etree.SubElement(sf2, qn('a:srgbClr')); s2.set('val', '0682b2')
    a2 = etree.SubElement(s2, qn('a:alpha')); a2.set('val', '20000')
c2.line.fill.background()

# Teal accent bar left
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=TEAL)

# Ministry hierarchy — top
add_textbox(slide, "جمهورية العراق  |  وزارة الصحة  |  دائرة صحة بابل",
            1, 0.25, 11.5, 0.5, font_size=12, color=RGBColor(0xa7,0xf3,0xd0),
            align=PP_ALIGN.CENTER, rtl=True)
add_textbox(slide, "قسم الصيدلة  —  شعبة الصيدلة السريرية",
            1, 0.65, 11.5, 0.4, font_size=11, color=RGBColor(0x6e,0xe7,0xb7),
            align=PP_ALIGN.CENTER, rtl=True)

# Main title Arabic
add_textbox(slide, "نظام تحليل حساسية المضادات الحيوية",
            0.5, 1.5, 12.3, 1.1, font_size=40, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, rtl=True)

# Subtitle Arabic
add_textbox(slide, "المرجع الإكلينيكي المتكامل لرصد المقاومة الميكروبية وتحليل الأنماط الوبائية",
            0.5, 2.55, 12.3, 0.65, font_size=17, color=RGBColor(0xa7,0xf3,0xd0),
            align=PP_ALIGN.CENTER, rtl=True)

# Divider
add_rect(slide, 2.5, 3.3, 8.33, 0.04, fill_rgb=TEAL)

# English title
add_textbox(slide, "AntibioGram Pro",
            0.5, 3.45, 12.3, 0.75, font_size=34, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

# English subtitle
add_textbox(slide, "Integrated Clinical Platform for Antimicrobial Resistance Surveillance & Antibiogram Analysis",
            0.5, 4.15, 12.3, 0.55, font_size=13,
            color=RGBColor(0x99,0xf6,0xe4), align=PP_ALIGN.CENTER)

# Standards badges
for i, (label, x) in enumerate([("CLSI M100", 3.8), ("EUCAST 2024", 5.55), ("Magiorakos 2012", 7.3)]):
    add_rect(slide, x, 4.88, 1.6, 0.38, fill_rgb=RGBColor(0x13,0x4e,0x4a),
             line_rgb=TEAL, line_width=1)
    add_textbox(slide, label, x + 0.05, 4.9, 1.5, 0.32,
                font_size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Bottom attribution
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=RGBColor(0x06,0x24,0x22))
add_textbox(slide,
            "وحدة متابعة لجان الصيدلة والعلاج  |  Abdallah Jawad Kadhim  |  © 2026",
            0.3, 6.87, 12.7, 0.5, font_size=12,
            color=RGBColor(0x5e,0xea,0xd4), align=PP_ALIGN.CENTER, rtl=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "نظرة عامة على المنظومة", "System Overview", y=0.18)

cards = [
    ("📊", "لوحة البيانات", "رصد فوري للمقاومة\nومؤشرات الجودة",   TEAL,    1.0),
    ("🏥", "إدارة المستشفيات", "دعم شبكة مستشفيات\nمتعددة المواقع", CYAN,    3.6),
    ("🧬", "تحليل الأنماط",    "MRSA · VRE · CRE\nMDR · XDR · PDR", VIOLET,  6.2),
    ("📈", "الاتجاهات",        "توقعات الانحدار\nOLS الخطي",        EMERALD, 8.8),
    ("📋", "التقارير",         "تصدير Excel/PDF\nأو Word",          AMBER,  11.4),
]

for icon, title_ar, desc_ar, clr, x in cards:
    add_rect(slide, x, 1.1, 1.8, 5.9, fill_rgb=RGBColor(0xf8,0xff,0xfe),
             line_rgb=clr, line_width=1.5)
    add_rect(slide, x, 1.1, 1.8, 0.55, fill_rgb=clr)
    add_textbox(slide, icon,     x, 1.12, 1.8, 0.48, font_size=20, align=PP_ALIGN.CENTER)
    add_textbox(slide, title_ar, x, 1.7,  1.8, 0.55,
                font_size=12, bold=True, color=clr,
                align=PP_ALIGN.CENTER, rtl=True)
    add_textbox(slide, desc_ar,  x, 2.3,  1.8, 1.3,
                font_size=10, color=MID, align=PP_ALIGN.CENTER, rtl=True)

# Bottom note
add_rect(slide, 0.4, 7.05, 12.53, 0.35, fill_rgb=RGBColor(0xec,0xfd,0xfb),
         line_rgb=TEAL, line_width=1)
add_textbox(slide,
            "يعمل البرنامج محلياً بالكامل — لا اتصال إنترنت مطلوب — البيانات مشفّرة ومحمية",
            0.5, 7.06, 12.4, 0.3, font_size=11, color=TEAL,
            align=PP_ALIGN.CENTER, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Mathematical Core: SIR Formulas
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "المنطق الرياضي الأساسي — معادلات SIR", "Core Mathematics — SIR Formulas")

add_textbox(slide,
            "يُحسب كل معامل بشكل مستقل من العدد الكلي N وليس بالتكملة — هذا يضمن الدقة عند وجود فئة الحساسية المتوسطة (I)",
            0.5, 1.1, 12.3, 0.45, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

# Three formula boxes
add_formula_box(slide, "%S  — Susceptible",
                "%S = (S ÷ N) × 100",
                "S = عدد المعزولات الحساسة    N = العدد الكلي", 0.8, 1.7)

add_formula_box(slide, "%I  — Intermediate",
                "%I = (I ÷ N) × 100",
                "I = عدد المعزولات متوسطة الحساسية", 4.65, 1.7)

add_formula_box(slide, "%R  — Resistant",
                "%R = (R ÷ N) × 100",
                "الخطأ الشائع: %R ≠ 100 − %S  عند وجود I > 0", 8.5, 1.7)

# Why section
add_rect(slide, 0.8, 3.05, 11.73, 0.55, fill_rgb=RGBColor(0xff,0xf7,0xed),
         line_rgb=AMBER, line_width=1.2)
add_textbox(slide,
            "⚠  لماذا لا يُستخدم %R = 100 − %S ؟  لأن هذه الصيغة تُسقط الفئة المتوسطة (I) وتُبالغ في تقدير المقاومة بشكل خاطئ إكلينيكياً",
            0.9, 3.07, 11.5, 0.45, font_size=11, color=AMBER,
            align=PP_ALIGN.RIGHT, rtl=True)

# Reliability threshold
add_rect(slide, 0.8, 3.75, 11.73, 1.05, fill_rgb=RGBColor(0xf0,0xfd,0xfa),
         line_rgb=TEAL, line_width=1.2)
add_textbox(slide, "معيار الموثوقية — CLSI M39", 0.9, 3.8, 11.5, 0.38,
            font_size=14, bold=True, color=TEAL, align=PP_ALIGN.RIGHT, rtl=True)
add_textbox(slide,
            "N ≥ 30 معزولة  →  النتيجة موثوقة إحصائياً\n"
            "N < 30 معزولة  →  تُعرض النتيجة مع تحذير (*) وفق معيار CLSI M39-A4",
            0.9, 4.18, 11.5, 0.55, font_size=12, color=DARK,
            align=PP_ALIGN.RIGHT, rtl=True)

# Test case table
add_textbox(slide, "مثال تطبيقي:", 0.8, 4.97, 3, 0.3,
            font_size=12, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

headers = ["S", "I", "R", "N", "%S", "%I", "%R", "موثوق"]
vals    = ["15", "5", "10", "30", "50.0%", "16.7%", "33.3%", "✓"]
col_w   = 1.45
for j, (h, v) in enumerate(zip(headers, vals)):
    cx = 0.8 + j * col_w
    add_rect(slide, cx, 5.28, col_w, 0.3, fill_rgb=TEAL)
    add_textbox(slide, h, cx, 5.3, col_w, 0.24, font_size=10,
                bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, cx, 5.58, col_w, 0.3, fill_rgb=LIGHT_BG,
             line_rgb=TEAL, line_width=0.5)
    add_textbox(slide, v, cx, 5.6, col_w, 0.24, font_size=10,
                color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CLSI / EUCAST Breakpoints
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "نقاط الانكسار — معايير CLSI و EUCAST", "Antimicrobial Breakpoints — CLSI & EUCAST")

add_textbox(slide,
            "يستخدم البرنامج جداول نقاط الانكسار المُحدَّثة لتصنيف المعزولات إكلينيكياً وفق تركيز MIC أو قطر منطقة التثبيط",
            0.5, 1.1, 12.3, 0.4, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

# Two-column comparison
for i, (std, clr, x, items) in enumerate([
    ("CLSI  M100", TEAL, 0.4, [
        "نقطتا الانكسار:  S ≤ X  |  I بين X و Y  |  R ≥ Y",
        "منطقة الحساسية المتوسطة (I) = قابلة للعلاج بجرعات مرتفعة",
        "تُطبَّق على الجرثومة المُعزولة في المختبر السريري",
        "يشمل: Enterobacterales · Pseudomonas · Acinetobacter",
        "Staphylococcus · Enterococcus · Streptococcus",
    ]),
    ("EUCAST  2024", VIOLET, 6.9, [
        "نقطتا الانكسار:  S ≤ X  |  R > X  (لا توجد فئة I دائماً)",
        'اصطلاح "R > X": المعزولة بـ MIC = X تُصنَّف متوسطة وليس مقاومة',
        "rLimit في البرنامج = 2× نقطة الانكسار لـ R (التخفيف المضاعف التالي)",
        "يشمل نفس مجموعات الجراثيم مع فوارق في القيم العددية",
        "يُطبَّق عند اختيار المعيار الأوروبي في إعدادات البرنامج",
    ]),
]):
    add_rect(slide, x, 1.6, 6.0, 5.3, fill_rgb=RGBColor(0xfa,0xff,0xfe),
             line_rgb=clr, line_width=2)
    add_rect(slide, x, 1.6, 6.0, 0.48, fill_rgb=clr)
    add_textbox(slide, std, x, 1.63, 6.0, 0.4,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.18),
                                   Inches(5.7), Inches(4.5))
    tf = txb.text_frame; tf.word_wrap = True
    for item in items:
        add_paragraph(tf, "  •  " + item, font_size=12, color=DARK,
                      align=PP_ALIGN.RIGHT, rtl=True, space_before=6)

# Bottom note on EUCAST rLimit rule
add_rect(slide, 0.4, 7.05, 12.53, 0.38, fill_rgb=RGBColor(0xed,0xe9,0xfe),
         line_rgb=VIOLET, line_width=1)
add_textbox(slide,
            "قاعدة EUCAST الحاسمة: للتراكيز غير المتصلة — rLimit = نقطة R المنشورة × 2  (التخفيف المضاعف التالي)",
            0.5, 7.08, 12.4, 0.3, font_size=11, color=VIOLET,
            align=PP_ALIGN.RIGHT, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Wilson 95% Confidence Interval
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "فترة الثقة ويلسون 95%", "Wilson 95% Confidence Interval")

add_textbox(slide,
            "تُقدِّم فترة الثقة المدى الحقيقي المتوقع للحساسية في المجتمع الميكروبي — وهي أفضل من فترة والد-بيرسون عند N صغيرة",
            0.5, 1.1, 12.3, 0.42, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

# Formula display
add_rect(slide, 1.5, 1.65, 10.33, 1.6, fill_rgb=RGBColor(0xf0,0xfd,0xfa),
         line_rgb=TEAL, line_width=1.5)
add_textbox(slide, "صيغة ويلسون:", 1.6, 1.7, 4, 0.35,
            font_size=12, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
add_textbox(slide,
            "p̂  =  S/N          z  =  1.959964  (95% CI)",
            2.0, 2.0, 9.5, 0.35, font_size=14, color=NAVY,
            align=PP_ALIGN.CENTER, font_name="Cambria Math")
add_textbox(slide,
            "CI  =  ( p̂ + z²/2N  ±  z√[p̂(1−p̂)/N + z²/4N²] ) / ( 1 + z²/N )",
            1.6, 2.38, 10.0, 0.45, font_size=14, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, font_name="Cambria Math")

# Why Wilson
add_textbox(slide, "لماذا ويلسون وليس والد-بيرسون؟", 0.6, 3.42, 12, 0.38,
            font_size=14, bold=True, color=TEAL, align=PP_ALIGN.RIGHT, rtl=True)

rows_why = [
    ("ملاحظة", "فترة Wald", "فترة Wilson"),
    ("N صغيرة (< 30)",     "قد تتجاوز [0,1]",   "مقيدة دائماً في [0,1]"),
    ("p̂ = 0 أو 1",         "CI = [0,0] أو [1,1]","CI واسع وصحيح إحصائياً"),
    ("الموثوقية العامة",    "ضعيفة عند الأطراف", "ممتازة — توصية CLSI"),
]
col_ws = [3.5, 4.0, 4.0]; xs = [0.6, 4.1, 8.1]
for ri, row in enumerate(rows_why):
    for ci, (cell, cw, cx) in enumerate(zip(row, col_ws, xs)):
        bg = TEAL if ri == 0 else (RGBColor(0xf0,0xfd,0xfa) if ri % 2 == 0 else WHITE)
        fc = WHITE if ri == 0 else DARK
        add_rect(slide, cx, 3.85 + ri*0.42, cw, 0.42,
                 fill_rgb=bg, line_rgb=TEAL, line_width=0.5)
        add_textbox(slide, cell, cx + 0.05, 3.87 + ri*0.42, cw - 0.1, 0.36,
                    font_size=11, bold=(ri==0), color=fc, align=PP_ALIGN.CENTER, rtl=(ci==0))

# Example
add_textbox(slide,
            "مثال: S=15, I=5, R=10, N=30  →  %S = 50.0%  |  95% CI  [31.9% — 68.1%]",
            0.6, 5.6, 12.2, 0.38, font_size=12, color=EMERALD,
            align=PP_ALIGN.CENTER, rtl=False)
add_textbox(slide,
            "يعني: الحساسية الفعلية في المجتمع تقع بين 31.9% و 68.1% بثقة 95%",
            0.6, 5.97, 12.2, 0.35, font_size=11, color=MID,
            align=PP_ALIGN.CENTER, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AMR Phenotype Detection
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "كشف الأنماط المقاومة — AMR Phenotypes", "AMR Phenotype Detection")

add_textbox(slide,
            "يرصد البرنامج تلقائياً الأنماط الخطيرة بناءً على قواعد الارتباط بين الجرثومة والمضاد المُبلَّغ عنه",
            0.5, 1.1, 12.3, 0.4, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

phenotypes = [
    ("MRSA", "المكوّر العنقودي الذهبي المقاوم للميثيسيلين",
     "S. aureus + Oxacillin / Cefoxitin مقاومة", ROSE, "خطر اكتساب العدوى وصعوبة العلاج"),
    ("VRE",  "المكوّر المعوي المقاوم للفانكومايسين",
     "Enterococcus spp. + Vancomycin مقاومة",   VIOLET, "خيارات علاجية محدودة للغاية"),
    ("CRE",  "الجراثيم المعوية المقاومة للكاربابينيم",
     "Enterobacterales + Meropenem/Imipenem/Ertapenem", AMBER, "أزمة صحية عالمية — WHO Priority 1"),
    ("ESBL", "مُنتِجات بيتا-لاكتاماز الطيف الواسع",
     "Enterobacterales + Ceftriaxone/Ceftazidime/Cefepime", EMERALD, "تفشي واسع في العيادات"),
    ("CRAB", "A. baumannii المقاومة للكاربابينيم",
     "Acinetobacter baumannii + Carbapenems",    CYAN, "شائع في وحدات العناية المركزة"),
    ("CRPA", "P. aeruginosa المقاومة للكاربابينيم",
     "Pseudomonas aeruginosa + Meropenem/Imipenem", RGBColor(0x9f,0x1a,0xf5), "أولوية حرجة — WHO"),
]

for i, (code, name_ar, trigger, clr, impact) in enumerate(phenotypes):
    r, c = divmod(i, 3)
    x = 0.4 + c * 4.3
    y = 1.65 + r * 2.65
    add_rect(slide, x, y, 4.1, 2.45, fill_rgb=WHITE, line_rgb=clr, line_width=2)
    add_rect(slide, x, y, 4.1, 0.48, fill_rgb=clr)
    add_textbox(slide, code, x, y + 0.04, 4.1, 0.38,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, name_ar, x + 0.1, y + 0.52, 3.9, 0.42,
                font_size=10, bold=True, color=clr, align=PP_ALIGN.CENTER, rtl=True)
    add_textbox(slide, trigger, x + 0.1, y + 0.94, 3.9, 0.42,
                font_size=9, color=MID, align=PP_ALIGN.CENTER, rtl=False)
    add_rect(slide, x + 0.1, y + 1.45, 3.9, 0.04, fill_rgb=RGBColor(0xe2,0xe8,0xf0))
    add_textbox(slide, "⚑ " + impact, x + 0.1, y + 1.54, 3.9, 0.38,
                font_size=9, color=ROSE, align=PP_ALIGN.RIGHT, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MDR / XDR / PDR Classification
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "تصنيف المقاومة المتعددة  MDR · XDR · PDR",
                   "Multi-Drug Resistance Classification  (Magiorakos et al. 2012, CMI)")

add_textbox(slide,
            'تعريفات إجماعية دولية نُشرت عام 2012 في مجلة Clinical Microbiology & Infection — "تعريفات موحدة للجراثيم المعزولة المكتسبة للمقاومة"',
            0.5, 1.12, 12.3, 0.42, font_size=11, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

levels = [
    ("MDR", "Multi-Drug Resistant",
     "مقاوم لـ ≥ 1 عامل في ≥ 3 فئات دوائية مختبَرة",
     "resistant_classes ≥ 3", AMBER, 0.55),
    ("XDR", "Extensively Drug Resistant",
     "مقاوم لجميع الفئات ما عدا ≤ 2 فئتين",
     "resistant ≥ total − 2  AND  resistant ≥ 3", ROSE, 4.55),
    ("PDR", "Pan-Drug Resistant",
     "مقاوم لجميع العوامل في جميع الفئات المختبَرة",
     "resistant_classes = total_classes", RGBColor(0x7f,0x1d,0x1d), 8.55),
]

for code, full, ar_def, code_logic, clr, x in levels:
    add_rect(slide, x, 1.7, 4.2, 5.0, fill_rgb=WHITE, line_rgb=clr, line_width=2)
    add_rect(slide, x, 1.7, 4.2, 0.6, fill_rgb=clr)
    add_textbox(slide, code, x, 1.73, 4.2, 0.5,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, full, x + 0.1, 2.37, 4.0, 0.4,
                font_size=11, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(slide, ar_def, x + 0.1, 2.82, 4.0, 0.55,
                font_size=12, bold=True, color=DARK, align=PP_ALIGN.RIGHT, rtl=True)
    add_rect(slide, x + 0.15, 3.5, 3.9, 0.04, fill_rgb=RGBColor(0xe2,0xe8,0xf0))
    add_textbox(slide, "منطق الكود:", x + 0.15, 3.58, 3.9, 0.3,
                font_size=9, color=MID, align=PP_ALIGN.LEFT)
    add_textbox(slide, code_logic, x + 0.15, 3.9, 3.9, 0.42,
                font_size=10, bold=True, color=clr, align=PP_ALIGN.LEFT,
                font_name="Consolas")

# 5% threshold note
add_rect(slide, 0.4, 6.82, 12.53, 0.55, fill_rgb=RGBColor(0xff,0xf7,0xed),
         line_rgb=AMBER, line_width=1.2)
add_textbox(slide,
            "عتبة 5%: في الأنتيبيوغرام التراكمي، تُعدّ الفئة الدوائية مقاومة فقط إذا بلغت نسبة المقاومة ≥ 5%\n"
            "هذا يمنع تضخيم التصنيف بسبب معزولة واحدة شاذة في عينة كبيرة.",
            0.5, 6.85, 12.3, 0.48, font_size=11, color=AMBER,
            align=PP_ALIGN.RIGHT, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WISCA Empiric Coverage
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "تغطية المضادات التجريبية — نموذج WISCA", "Empiric Antibiotic Coverage — WISCA Model")

add_textbox(slide,
            "يُحسب البرنامج احتمال أن يُغطّي مضاد مُختار كل الجراثيم المُسبِّبة للعدوى موزونةً بنسبة انتشارها",
            0.5, 1.1, 12.3, 0.4, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

# Formula
add_rect(slide, 1.0, 1.62, 11.33, 1.2, fill_rgb=RGBColor(0xf0,0xfd,0xfa),
         line_rgb=TEAL, line_width=1.5)
add_textbox(slide, "صيغة التغطية الموزونة:", 1.1, 1.67, 5, 0.35,
            font_size=13, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
add_textbox(slide,
            "Coverage(Ab) = Σ  [ P(organism_i)  ×  %S(organism_i, Ab) ]  /  Σ P(organism_i)",
            1.1, 2.0, 11.1, 0.45, font_size=14, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, font_name="Cambria Math")

# Bands
add_textbox(slide, "مستويات التغطية المُعتمدة إكلينيكياً:", 0.6, 2.97, 12, 0.38,
            font_size=13, bold=True, color=TEAL, align=PP_ALIGN.RIGHT, rtl=True)

bands = [
    ("≥ 90%", "ممتازة", "Excellent", EMERALD),
    ("80–89%", "جيدة", "Good", TEAL),
    ("70–79%", "مقبولة مع حذر", "Acceptable", AMBER),
    ("< 70%", "غير كافية", "Insufficient", ROSE),
]
for i, (pct, ar, en, clr) in enumerate(bands):
    x = 0.5 + i * 3.08
    add_rect(slide, x, 3.4, 3.0, 1.1, fill_rgb=WHITE, line_rgb=clr, line_width=2)
    add_rect(slide, x, 3.4, 3.0, 0.38, fill_rgb=clr)
    add_textbox(slide, pct, x, 3.42, 3.0, 0.32,
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, ar, x, 3.82, 3.0, 0.35,
                font_size=13, bold=True, color=clr, align=PP_ALIGN.CENTER, rtl=True)
    add_textbox(slide, en, x, 4.2, 3.0, 0.28,
                font_size=10, color=MID, align=PP_ALIGN.CENTER)

# Clinical use case
add_rect(slide, 0.5, 4.68, 12.33, 1.82, fill_rgb=RGBColor(0xf8,0xfa,0xff),
         line_rgb=VIOLET, line_width=1.2)
add_textbox(slide, "حالة سريرية كمثال:", 0.65, 4.73, 5, 0.35,
            font_size=13, bold=True, color=VIOLET, align=PP_ALIGN.LEFT)
example_lines = [
    "مريض مشتبه بعدوى بولية — توزيع الجراثيم المحلية: E.coli 60% | K.pneumoniae 25% | P.aeruginosa 15%",
    "مضاد A (Ceftriaxone): %S = 70%، 55%، 30%  →  تغطية موزونة = 0.6×70 + 0.25×55 + 0.15×30 = 60.25% ✗",
    "مضاد B (Piperacillin-Tazobactam): %S = 85%، 80%، 78%  →  تغطية = 82.7%  ✓ جيدة",
]
txb = slide.shapes.add_textbox(Inches(0.65), Inches(5.12), Inches(12.0), Inches(1.3))
tf = txb.text_frame; tf.word_wrap = True
for line in example_lines:
    add_paragraph(tf, line, font_size=10, color=DARK,
                  align=PP_ALIGN.RIGHT, rtl=True, space_before=4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Statistical Analysis
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "التحليل الإحصائي — مقارنة الاتجاهات والتنبؤ", "Statistical Analysis — Trend Comparison & Forecasting")

# Chi-square
add_rect(slide, 0.4, 1.1, 5.9, 5.8, fill_rgb=RGBColor(0xf0,0xfd,0xfa),
         line_rgb=TEAL, line_width=1.5)
add_rect(slide, 0.4, 1.1, 5.9, 0.5, fill_rgb=TEAL)
add_textbox(slide, "اختبار مربع كاي (χ²)", 0.5, 1.13, 5.7, 0.4,
            font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)
txb = slide.shapes.add_textbox(Inches(0.55), Inches(1.72), Inches(5.6), Inches(5.0))
tf = txb.text_frame; tf.word_wrap = True
chi_points = [
    "مقارنة معدلات المقاومة بين مستشفيين أو فترتين زمنيتين",
    "جدول 2×2: مقاوم / غير مقاوم × المجموعة أ/ب",
    "تصحيح ياتس Yates عند التوقعات < 5",
    "χ² = N(|ad−bc| − N/2)² / (r₁r₂c₁c₂)",
    "p-value من تقريب التوزيع الطبيعي Z = √χ²",
    "الدلالة الإحصائية: p < 0.05  (ثقة 95%)",
]
for pt in chi_points:
    add_paragraph(tf, "  •  " + pt, font_size=11, color=DARK,
                  align=PP_ALIGN.RIGHT, rtl=True, space_before=7)

# OLS Regression
add_rect(slide, 6.9, 1.1, 6.03, 5.8, fill_rgb=RGBColor(0xf5,0xf3,0xff),
         line_rgb=VIOLET, line_width=1.5)
add_rect(slide, 6.9, 1.1, 6.03, 0.5, fill_rgb=VIOLET)
add_textbox(slide, "انحدار OLS الخطي — التنبؤ بالمقاومة", 7.0, 1.13, 5.8, 0.4,
            font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)
txb2 = slide.shapes.add_textbox(Inches(7.05), Inches(1.72), Inches(5.7), Inches(3.5))
tf2 = txb2.text_frame; tf2.word_wrap = True
ols_points = [
    "تحليل معدل المقاومة عبر سنوات متعددة",
    "slope = Σ(xᵢ−x̄)(yᵢ−ȳ) / Σ(xᵢ−x̄)²",
    "intercept = ȳ − slope × x̄",
    "y = معدل المقاومة (100 − %S)",
    "التنبؤ: resistance_2026 = slope×2026 + intercept",
    "الحصر في [0%, 100%] تلقائياً",
]
for pt in ols_points:
    add_paragraph(tf2, "  •  " + pt, font_size=11, color=DARK,
                  align=PP_ALIGN.RIGHT, rtl=True, space_before=7)

# Trend example
add_rect(slide, 6.9, 4.65, 6.03, 2.25, fill_rgb=WHITE, line_rgb=VIOLET, line_width=1)
add_textbox(slide, "مثال: Ciprofloxacin vs E. coli", 7.0, 4.7, 5.8, 0.35,
            font_size=11, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
years =   ["2021", "2022", "2023", "2024", "2025", "2026*"]
res   =   ["42%",  "47%",  "51%",  "55%",  "58%",  "62%"]
for j, (yr, rv) in enumerate(zip(years, res)):
    cx = 6.92 + j * 0.97
    add_rect(slide, cx, 5.12, 0.92, 0.28,
             fill_rgb=VIOLET if yr == "2026*" else RGBColor(0xed,0xe9,0xfe))
    add_textbox(slide, yr, cx, 5.13, 0.92, 0.22, font_size=8,
                color=WHITE if yr == "2026*" else VIOLET, align=PP_ALIGN.CENTER)
    add_rect(slide, cx, 5.42, 0.92, 0.28, fill_rgb=WHITE, line_rgb=VIOLET, line_width=0.5)
    add_textbox(slide, rv, cx, 5.43, 0.92, 0.22, font_size=9,
                bold=(yr == "2026*"), color=ROSE if yr == "2026*" else DARK,
                align=PP_ALIGN.CENTER)
add_textbox(slide, "* توقع رياضي بالانحدار الخطي", 6.9, 5.78, 6.0, 0.28,
            font_size=9, color=MID, align=PP_ALIGN.CENTER, italic=True)

# Bottom note
add_rect(slide, 0.4, 7.0, 12.53, 0.4, fill_rgb=RGBColor(0xf8,0xfa,0xff),
         line_rgb=VIOLET, line_width=1)
add_textbox(slide,
            "تُعرض جميع المقارنات مع قيمة p وتفسيرها اللغوي التلقائي — يتطلب ≥ نقطتي بيانات للانحدار",
            0.5, 7.03, 12.3, 0.32, font_size=11, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Data Quality & First-Isolate Rule
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "جودة البيانات ومعايير CLSI M39", "Data Quality — CLSI M39 Standards")

add_textbox(slide,
            "وثيقة CLSI M39-A4 هي المعيار الذهبي لإعداد الأنتيبيوغرام التراكمي وتحديد شروط الإبلاغ الموثوق",
            0.5, 1.1, 12.3, 0.4, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

rules = [
    ("قاعدة المعزولة الأولى\nFirst-Isolate Rule — CLSI M39",
     "تُحتسب معزولة واحدة فقط لكل مريض لكل نوع جرثومي في الفترة الواحدة\n"
     "يمنع تضخيم النتائج من مرضى متكررين — يُطبَّق تلقائياً عند وجود Patient ID",
     TEAL),
    ("حد الموثوقية الإحصائية\nN ≥ 30 Isolates",
     "أقل من 30 معزولة يعني اتساعاً كبيراً في فترة الثقة وعدم استقرار النسب\n"
     "يضع البرنامج علامة (*) تلقائياً على كل نتيجة بـ N < 30 ويُنبّه المستخدم",
     CYAN),
    ("سلامة الإدخال\nInput Validation",
     "S + I + R = N  —  التحقق الرياضي من اتساق البيانات\n"
     "الكشف التلقائي عن القيم الغائبة والأنماط الشاذة وسوء التصنيف",
     VIOLET),
    ("إزالة التكرار\nDeduplication",
     "تُعرَّف النسخة المكررة بـ: Patient ID + Organism + Specimen + Antibiotic\n"
     "الأولوية للمعزولة الأولى زمنياً — الباقي يُزال مع إشعار العدد للمستخدم",
     EMERALD),
]

for i, (title, detail, clr) in enumerate(rules):
    r, c = divmod(i, 2)
    x = 0.4 + c * 6.47
    y = 1.65 + r * 2.65
    add_rect(slide, x, y, 6.1, 2.45, fill_rgb=WHITE, line_rgb=clr, line_width=1.5)
    add_rect(slide, x, y, 6.1, 0.12, fill_rgb=clr)
    add_textbox(slide, title, x + 0.12, y + 0.18, 5.86, 0.55,
                font_size=12, bold=True, color=clr, align=PP_ALIGN.RIGHT, rtl=True)
    txb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.82),
                                   Inches(5.8), Inches(1.55))
    tf = txb.text_frame; tf.word_wrap = True
    for line in detail.split('\n'):
        add_paragraph(tf, line, font_size=11, color=DARK,
                      align=PP_ALIGN.RIGHT, rtl=True, space_before=5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Multi-Hospital Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_section_header(slide, "إدارة الشبكة الاستشفائية — تنسيق إقليمي متكامل",
                   "Multi-Hospital Network — Regional Coordination")

add_textbox(slide,
            "يدير البرنامج شبكة غير محدودة من المستشفيات بتنسيق مركزي كامل مع دعم تنسيق البيانات الإقليمية",
            0.5, 1.1, 12.3, 0.4, font_size=12, color=MID,
            align=PP_ALIGN.RIGHT, rtl=True)

# Architecture diagram (simplified)
# Central node
add_rect(slide, 5.5, 2.0, 2.33, 0.9, fill_rgb=TEAL)
add_textbox(slide, "قاعدة البيانات\nالمركزية", 5.5, 2.0, 2.33, 0.85,
            font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)

hosp_nodes = [
    ("مستشفى بابل\nالتعليمي", 0.3, 1.6, CYAN),
    ("مستشفى الحلة\nالعام",    0.3, 3.5, EMERALD),
    ("مستشفى\nالنساء",        3.2, 0.85, VIOLET),
    ("مستشفى\nالأطفال",       3.2, 3.95, AMBER),
    ("مستشفى\nالقاسم",        10.2, 1.6, ROSE),
    ("مستشفى\nالمنصورية",     10.2, 3.5, RGBColor(0x06,0x82,0xa0)),
]

for label, x, y, clr in hosp_nodes:
    add_rect(slide, x, y, 2.5, 0.78, fill_rgb=WHITE, line_rgb=clr, line_width=1.5)
    add_textbox(slide, label, x, y + 0.04, 2.5, 0.68,
                font_size=10, bold=True, color=clr, align=PP_ALIGN.CENTER, rtl=True)

# Babil Format box
add_rect(slide, 3.2, 2.2, 2.2, 0.78, fill_rgb=RGBColor(0xf5,0xf3,0xff),
         line_rgb=VIOLET, line_width=1.5)
add_textbox(slide, "تنسيق بابل\nBabil Format", 3.2, 2.24, 2.2, 0.68,
            font_size=10, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

add_rect(slide, 7.93, 2.2, 2.2, 0.78, fill_rgb=RGBColor(0xf0,0xfd,0xfa),
         line_rgb=TEAL, line_width=1.5)
add_textbox(slide, "اللوحة الإقليمية\nRegional Dashboard", 7.93, 2.24, 2.2, 0.68,
            font_size=9, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Feature bullets
add_textbox(slide, "مميزات التنسيق الإقليمي:", 0.5, 5.1, 6.0, 0.35,
            font_size=13, bold=True, color=TEAL, align=PP_ALIGN.RIGHT, rtl=True)
feats_l = [
    "ملف Excel واحد متعدد الأوراق (كل ورقة = مستشفى)",
    "كشف تلقائي لتنسيق بابل وإنشاء المستشفيات",
    "حذف جماعي للمستشفيات بتحديد متعدد",
]
feats_r = [
    "مقارنة معدلات المقاومة بين المستشفيات",
    "خريطة حرارية لمستوى الخطر الإقليمي",
    "تقارير موحدة للجهات الصحية العليا",
]
txb_l = slide.shapes.add_textbox(Inches(6.5), Inches(5.55), Inches(6.3), Inches(1.8))
tf_l = txb_l.text_frame; tf_l.word_wrap = True
txb_r = slide.shapes.add_textbox(Inches(0.5), Inches(5.55), Inches(6.3), Inches(1.8))
tf_r = txb_r.text_frame; tf_r.word_wrap = True
for f in feats_l:
    add_paragraph(tf_l, "  ✦  " + f, font_size=11, color=DARK,
                  align=PP_ALIGN.RIGHT, rtl=True, space_before=5)
for f in feats_r:
    add_paragraph(tf_r, "  ✦  " + f, font_size=11, color=DARK,
                  align=PP_ALIGN.RIGHT, rtl=True, space_before=5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Closing / Attribution
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
gradient_bg(slide, NAVY, RGBColor(0x06,0x4e,0x4e))
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=TEAL)

# Stars / decorative dots
for xi, yi in [(11.5,0.8),(12.2,1.6),(10.8,2.1),(12.5,0.3),(11.1,1.2)]:
    c = slide.shapes.add_shape(9, Inches(xi), Inches(yi), Inches(0.1), Inches(0.1))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE
    spx = c._element.spPr
    sfx = spx.find('.//' + qn('a:solidFill'))
    if sfx is not None:
        sx = sfx.find(qn('a:srgbClr'))
        if sx is None:
            sx = etree.SubElement(sfx, qn('a:srgbClr')); sx.set('val', 'ffffff')
        ax = etree.SubElement(sx, qn('a:alpha')); ax.set('val', '40000')
    c.line.fill.background()

add_textbox(slide, "شكراً لاهتمامكم", 0.5, 1.3, 12.3, 0.85,
            font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)
add_textbox(slide, "Thank You", 0.5, 2.1, 12.3, 0.55,
            font_size=22, color=RGBColor(0xa7,0xf3,0xd0), align=PP_ALIGN.CENTER)

add_rect(slide, 2.5, 2.85, 8.33, 0.04, fill_rgb=TEAL)

# Attribution hierarchy
hierarchy = [
    ("جمهورية العراق",                          16, True,  WHITE),
    ("وزارة الصحة",                              18, True,  RGBColor(0xa7,0xf3,0xd0)),
    ("دائرة صحة بابل",                           15, False, RGBColor(0x6e,0xe7,0xb7)),
    ("قسم الصيدلة",                              14, False, RGBColor(0x6e,0xe7,0xb7)),
    ("شعبة الصيدلة السريرية",                    13, False, RGBColor(0x5e,0xea,0xd4)),
    ("بالتعاون مع",                              11, False, RGBColor(0x99,0xf6,0xe4)),
    ("وحدة متابعة لجان الصيدلة والعلاج",        12, True,  RGBColor(0xa7,0xf3,0xd0)),
]
y = 3.02
for text, fs, bd, clr in hierarchy:
    add_textbox(slide, text, 0.5, y, 12.3, 0.38 if fs >= 14 else 0.3,
                font_size=fs, bold=bd, color=clr, align=PP_ALIGN.CENTER, rtl=True)
    y += 0.38 if fs >= 14 else 0.32

add_rect(slide, 3.5, y + 0.05, 6.33, 0.04, fill_rgb=RGBColor(0x2d,0x6a,0x65))
add_textbox(slide, "Abdallah Jawad Kadhim", 0.5, y + 0.18, 12.3, 0.45,
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "© 2026  —  AntibioGram Pro", 0.5, y + 0.62, 12.3, 0.35,
            font_size=13, color=RGBColor(0x5e,0xea,0xd4), align=PP_ALIGN.CENTER)

# Save
out_path = r"C:\Users\abdal\OneDrive\Desktop\AntibioGram_Pro_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
